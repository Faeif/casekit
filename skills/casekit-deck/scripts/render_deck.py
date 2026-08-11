#!/usr/bin/env python3
"""Render a CaseKit deck specification to an editable PowerPoint file."""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_COLORS = {
    "navy": "102A43", "blue": "1677FF", "teal": "0F9D8A", "amber": "F59E0B",
    "red": "DC2626", "light": "F5F7FA", "ink": "172B4D", "muted": "627D98",
}


def rgb(value):
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


def add_box(slide, x, y, w, h, fill, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_text(slide, text, x, y, w, h, *, size=18, color="172B4D", bold=False,
             font="Arial", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, items, x, y, w, h, *, font, color, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for index, item in enumerate(items or []):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"{index + 1}. {item}"
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(10)
    return box


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def fmt(value):
    return f"{value:,.1f}".rstrip("0").rstrip(".") if isinstance(value, float) else f"{value:,}" if isinstance(value, int) else str(value)


def render(spec):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    colors = {**DEFAULT_COLORS, **spec.get("theme", {})}
    meta = spec.get("meta", {})
    font_head = meta.get("font_head", "Arial")
    font_body = meta.get("font_body", "Arial")

    for number, item in enumerate(spec["slides"], 1):
        slide = prs.slides.add_slide(blank)
        set_background(slide, "FFFFFF")
        slide_type = item.get("type", "content")
        if slide_type == "cover":
            set_background(slide, colors["navy"])
            add_box(slide, 0.8, 1.05, 0.75, 0.09, colors["teal"])
            add_text(slide, item["headline"], 0.8, 1.35, 11.3, 2.1, size=34, color="FFFFFF", bold=True, font=font_head, valign=MSO_ANCHOR.TOP)
            add_text(slide, item.get("subhead", meta.get("subtitle", "")), 0.82, 3.65, 10.8, 0.8, size=18, color="D9E2EC", font=font_body)
            add_text(slide, meta.get("team", "CaseKit"), 0.82, 6.55, 5, 0.3, size=12, color="9FB3C8", bold=True, font=font_body)
        else:
            add_box(slide, 0, 0, 0.12, 7.5, colors["blue"])
            add_text(slide, item["headline"], 0.65, 0.42, 11.9, 0.78, size=25, color=colors["navy"], bold=True, font=font_head, valign=MSO_ANCHOR.TOP)
            if slide_type == "metric":
                add_box(slide, 0.7, 1.45, 4.15, 4.75, colors["light"], "D8E2EC", rounded=True)
                add_text(slide, item.get("metric", "—"), 1.05, 2.05, 3.45, 1.25, size=42, color=colors["blue"], bold=True, font=font_head, align=PP_ALIGN.CENTER)
                add_text(slide, item.get("label", ""), 1.05, 3.25, 3.45, 0.65, size=17, color=colors["navy"], bold=True, font=font_body, align=PP_ALIGN.CENTER)
                add_text(slide, item.get("comparison", ""), 1.05, 4.05, 3.45, 0.55, size=15, color=colors["teal"], font=font_body, align=PP_ALIGN.CENTER)
                add_bullets(slide, item.get("body", []), 5.45, 1.7, 6.85, 4.35, font=font_body, color=colors["ink"])
            elif slide_type == "funnel":
                stages = item.get("stages", [])
                maximum = max([float(stage.get("value", 0)) for stage in stages] or [1])
                row_h = min(0.9, 4.8 / max(len(stages), 1))
                for index, stage in enumerate(stages):
                    width = 7.2 * max(float(stage.get("value", 0)) / maximum, 0.12)
                    x = 6.2 - width / 2
                    y = 1.45 + index * row_h
                    add_box(slide, x, y, width, row_h - 0.08, colors["teal"] if index == len(stages) - 1 else colors["blue"])
                    add_text(slide, stage.get("label", ""), 9.95, y, 1.35, row_h - 0.08, size=14, color=colors["navy"], font=font_body)
                    add_text(slide, fmt(stage.get("value", "")), 11.15, y, 1.1, row_h - 0.08, size=15, color=colors["ink"], bold=True, font=font_body, align=PP_ALIGN.RIGHT)
            elif slide_type == "timeline":
                phases = item.get("phases", [])
                width = 11.4 / max(len(phases), 1)
                for index, phase in enumerate(phases):
                    x = 0.75 + index * width
                    add_box(slide, x, 1.65, width - 0.2, 4.75, "EAF4FF" if index % 2 else colors["light"], "D8E2EC", rounded=True)
                    add_text(slide, phase.get("label", f"Phase {index + 1}"), x + 0.18, 1.9, width - 0.55, 0.55, size=18, color=colors["blue"], bold=True, font=font_body)
                    add_bullets(slide, phase.get("items", []), x + 0.15, 2.65, width - 0.5, 2.9, font=font_body, color=colors["ink"], size=14)
                    add_text(slide, phase.get("gate", ""), x + 0.18, 5.75, width - 0.55, 0.35, size=11, color=colors["teal"], bold=True, font=font_body)
            else:
                if item.get("ask"):
                    add_box(slide, 0.75, 1.55, 11.8, 1.25, colors["navy"], rounded=True)
                    add_text(slide, item["ask"], 1.05, 1.8, 11.2, 0.75, size=24, color="FFFFFF", bold=True, font=font_head, align=PP_ALIGN.CENTER)
                    add_bullets(slide, item.get("body", []), 1.2, 3.35, 10.8, 2.45, font=font_body, color=colors["ink"])
                else:
                    add_bullets(slide, item.get("body", []), 1.0, 1.65, 11.2, 4.9, font=font_body, color=colors["ink"])
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.08), Inches(12.25), Inches(0.01)).fill.solid()
            ids = item.get("evidence_ids", [])
            if ids:
                add_text(slide, "Evidence: " + " · ".join(ids), 0.65, 7.12, 11.5, 0.18, size=8.5, color=colors["muted"], font=font_body)
            add_text(slide, number, 12.5, 7.12, 0.3, 0.18, size=8.5, color=colors["muted"], font=font_body, align=PP_ALIGN.RIGHT)
        if item.get("speaker_notes"):
            try:
                slide.notes_slide.notes_text_frame.text = str(item["speaker_notes"])
            except AttributeError:
                pass
    return prs


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    spec = json.loads(args.spec.read_text(encoding="utf-8"))
    if not isinstance(spec.get("slides"), list) or not spec["slides"]:
        raise SystemExit("slides must be a non-empty array")
    for index, slide in enumerate(spec["slides"], 1):
        if not slide.get("headline"):
            raise SystemExit(f"slide {index} is missing headline")
    args.output.parent.mkdir(parents=True, exist_ok=True)
    render(spec).save(args.output)
    print(f"Rendered {len(spec['slides'])} slides -> {args.output}")


if __name__ == "__main__":
    main()
